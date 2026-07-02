# -*- coding: utf-8 -*-
import sys
from pptx import Presentation

def analizar_pptx():
    try:
        sys.stdout.reconfigure(encoding='utf-8')
    except AttributeError:
        pass
    pptx_path = r"c:\Users\ENRIQ\SISTEMA DE VOTANTES PARA EL ESTADO DE SONORA\capacitacion_brigadistas.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Error: No se encontró la presentación en {pptx_path}")
        return

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error al cargar la presentación: {e}")
        return

    print(f"=== PRESENTACIÓN: {os.path.basename(pptx_path)} ===")
    print(f"Total de Diapositivas: {len(prs.slides)}")
    print("=========================================\n")

    for i, slide in enumerate(prs.slides):
        print(f"--- Diapositiva {i+1} ---")
        
        # Intentar buscar un título o extraer todo el texto de las formas
        shapes_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    shapes_text.append(text)
            
            # Si tiene un grupo de formas o subformas
            elif shape.shape_type == 6: # Group shape
                for subshape in shape.shapes:
                    if subshape.has_text_frame:
                        text = subshape.text_frame.text.strip()
                        if text:
                            shapes_text.append(text)
                            
        if shapes_text:
            for text in shapes_text:
                print(f"[{text}]")
        else:
            print("[Diapositiva sin texto visible]")
        print("-" * 40 + "\n")

if __name__ == "__main__":
    import os
    analizar_pptx()
