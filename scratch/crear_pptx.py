# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# PALETA DE COLORES (Guinda & Oro)
# ==========================================
COLOR_GUINDA = RGBColor(166, 3, 33)      # Accent principal (#A60321)
COLOR_ORO = RGBColor(212, 175, 55)       # Oro brillante (#D4AF37)
COLOR_FONDO = RGBColor(249, 249, 249)    # Fondo gris claro (#F9F9F9)
COLOR_TEXTO_OSCURO = RGBColor(40, 40, 40)# Texto oscuro (#282828)
COLOR_BLANCO = RGBColor(255, 255, 255)   # Blanco
COLOR_GRIS_CLARO = RGBColor(230, 230, 230)

def aplicar_estilo_forma(forma, color_fondo, color_borde=None, grosor_borde=1):
    """Aplica color de fondo y borde a una forma de PPTX."""
    forma.fill.solid()
    forma.fill.fore_color.rgb = color_fondo
    if color_borde:
        forma.line.color.rgb = color_borde
        forma.line.width = Pt(grosor_borde)
    else:
        forma.line.fill.background()

def agregar_encabezado(slide, titulo_texto):
    """Agrega una barra superior estilizada y el título de la diapositiva."""
    # Barra superior oro
    barra_oro = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    aplicar_estilo_forma(barra_oro, COLOR_ORO)
    
    # Barra superior guinda
    barra_guinda = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.15), Inches(13.333), Inches(0.1))
    aplicar_estilo_forma(barra_guinda, COLOR_GUINDA)

    # Cuadro de título
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo_texto
    p.font.name = 'Georgia'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_GUINDA
    
    # Línea divisoria elegante
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.733), Inches(0.02))
    aplicar_estilo_forma(div, COLOR_GRIS_CLARO)

def crear_presentacion():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ------------------------------------------
    # Diapositiva 1: Portada
    # ------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_FONDO)

    # Panel lateral guinda
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4.0), Inches(7.5))
    aplicar_estilo_forma(panel, COLOR_GUINDA)

    # Franja oro decorativa vertical
    franja = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), 0, Inches(0.15), Inches(7.5))
    aplicar_estilo_forma(franja, COLOR_ORO)

    # Textos de Portada en el Panel Guinda
    txBox_panel = slide.shapes.add_textbox(Inches(0.4), Inches(2.2), Inches(3.2), Inches(3.5))
    tf_panel = txBox_panel.text_frame
    tf_panel.word_wrap = True
    
    p_dep = tf_panel.paragraphs[0]
    p_dep.text = "SISTEMA DE VOTANTES"
    p_dep.font.name = 'Arial'
    p_dep.font.size = Pt(14)
    p_dep.font.bold = True
    p_dep.font.color.rgb = COLOR_ORO
    p_dep.alignment = PP_ALIGN.LEFT
    
    p_camp = tf_panel.add_paragraph()
    p_camp.text = "Campaña Territorial\nSonora 2026"
    p_camp.font.name = 'Arial'
    p_camp.font.size = Pt(12)
    p_camp.font.color.rgb = COLOR_BLANCO
    p_camp.space_before = Pt(8)

    # Título Principal (Lado derecho)
    txBox_titulo = slide.shapes.add_textbox(Inches(4.8), Inches(2.0), Inches(7.8), Inches(4.0))
    tf_titulo = txBox_titulo.text_frame
    tf_titulo.word_wrap = True
    
    p_t = tf_titulo.paragraphs[0]
    p_t.text = "Manual de Capacitación para Brigadistas"
    p_t.font.name = 'Georgia'
    p_t.font.size = Pt(46)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_GUINDA
    
    p_sub = tf_titulo.add_paragraph()
    p_sub.text = "Sistema de Registro y Vinculación Territorial"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_TEXTO_OSCURO
    p_sub.space_before = Pt(14)

    p_p = tf_titulo.add_paragraph()
    p_p.text = "Senador Heriberto Aguilar Castillo"
    p_p.font.name = 'Georgia'
    p_p.font.size = Pt(16)
    p_p.font.italic = True
    p_p.font.bold = True
    p_p.font.color.rgb = COLOR_ORO
    p_p.space_before = Pt(20)

    # ------------------------------------------
    # Diapositiva 2: El Rol del Brigadista
    # ------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_FONDO)
    agregar_encabezado(slide, "El Rol del Brigadista en Territorio")

    # Contenido (Caja de Texto)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    bullets = [
        "Facilitar el registro organizado de ciudadanos y enlaces territoriales directos.",
        "Capturar datos demográficos reales de la zona de brigadeo para habilitar el cruce estadístico.",
        "Identificar y canalizar de forma inmediata peticiones de apoyo social de la ciudadanía.",
        "Promover la participación ciudadana y consolidar la estructura del movimiento."
    ]
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.space_before = Pt(14)

    # Tarjeta Destacada Lateral
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
    aplicar_estilo_forma(card, COLOR_GUINDA)
    
    txBox_card = slide.shapes.add_textbox(Inches(9.0), Inches(2.5), Inches(3.3), Inches(3.0))
    tf_card = txBox_card.text_frame
    tf_card.word_wrap = True
    p_card = tf_card.paragraphs[0]
    p_card.text = "OBJETIVO GENERAL\n\nConstruir una red organizada, confiable y con comunicación directa entre el Senador y la ciudadanía."
    p_card.font.name = 'Arial'
    p_card.font.size = Pt(16)
    p_card.font.bold = True
    p_card.font.color.rgb = COLOR_BLANCO
    p_card.alignment = PP_ALIGN.CENTER

    # ------------------------------------------
    # Diapositiva 3: Paso 1 - Enlace Único
    # ------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_FONDO)
    agregar_encabezado(slide, "Paso 1: Acceso con tu Enlace Único")

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = [
        "Tu coordinador te proporcionará una liga única de registro que contiene tu ID de brigadista.",
        "Muestra el código QR asociado a esta liga para que el ciudadano lo escanee con su propio celular.",
        "Si no puede escanearlo, compártele el enlace por WhatsApp para que lo abra en su dispositivo."
    ]
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.space_before = Pt(14)

    # Tarjeta de Advertencia (Warning Card)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
    aplicar_estilo_forma(card, COLOR_BLANCO, COLOR_GUINDA, 2)
    
    txBox_card = slide.shapes.add_textbox(Inches(8.9), Inches(2.2), Inches(3.5), Inches(3.8))
    tf_card = txBox_card.text_frame
    tf_card.word_wrap = True
    
    p_warn_t = tf_card.paragraphs[0]
    p_warn_t.text = "⚠️ ADVERTENCIA"
    p_warn_t.font.name = 'Arial'
    p_warn_t.font.size = Pt(18)
    p_warn_t.font.bold = True
    p_warn_t.font.color.rgb = COLOR_GUINDA
    p_warn_t.alignment = PP_ALIGN.CENTER
    
    p_warn_b = tf_card.add_paragraph()
    p_warn_b.text = "\nNUNCA registres al votante usando tu celular.\n\nLa sesión de WhatsApp de quien envía el mensaje final es la que queda asociada al registro. Si usas tu cel, los datos del votante se ligarán a tu número."
    p_warn_b.font.name = 'Arial'
    p_warn_b.font.size = Pt(13)
    p_warn_b.font.bold = True
    p_warn_b.font.color.rgb = COLOR_TEXTO_OSCURO
    p_warn_b.alignment = PP_ALIGN.CENTER

    # ------------------------------------------
    # Diapositiva 4: Paso 2 - Registro en la Web
    # ------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_FONDO)
    agregar_encabezado(slide, "Paso 2: Formulario Web del Ciudadano")

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = [
        "El votante entra a la liga y verá un banner arriba confirmando quién lo invita: 'Te invita el Enlace: [Tu Nombre]'.",
        "Ayúdale a capturar sus datos obligatorios: Nombre completo (como aparece en su INE), WhatsApp de 10 dígitos, Calle y Sección Electoral.",
        "Sección Electoral: Es el número de 3 o 4 dígitos ubicado al frente de la credencial INE.",
        "El votante puede elegir si desea sumarse al movimiento como Promotor, Defensor o Activista."
    ]
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.space_before = Pt(14)

    # Tarjeta Informativa Lateral
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
    aplicar_estilo_forma(card, COLOR_ORO)
    
    txBox_card = slide.shapes.add_textbox(Inches(9.0), Inches(2.5), Inches(3.3), Inches(3.0))
    tf_card = txBox_card.text_frame
    tf_card.word_wrap = True
    p_card = tf_card.paragraphs[0]
    p_card.text = "VÍNCULO DE CONFIANZA\n\nVerificar visualmente tu nombre en la parte superior del formulario le asegura al votante que está registrándose en el canal correcto contigo."
    p_card.font.name = 'Arial'
    p_card.font.size = Pt(15)
    p_card.font.bold = True
    p_card.font.color.rgb = COLOR_BLANCO
    p_card.alignment = PP_ALIGN.CENTER

    # ------------------------------------------
    # Diapositiva 5: DIAGRAMA DE FLUJO (FLOWCHART DIBUJADO)
    # ------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_FONDO)
    agregar_encabezado(slide, "Diagrama del Flujo de Trabajo")

    # Cajas del Flujo
    cajas_info = [
        {"x": 0.8, "y": 2.5, "w": 2.2, "h": 2.0, "t": "1. Escaneo QR\n\nCiudadano escanea la liga única de red del brigadista."},
        {"x": 3.8, "y": 2.5, "w": 2.2, "h": 2.0, "t": "2. Web Form\n\nSe abre la web, se confirma el líder y se capturan datos INE."},
        {"x": 6.8, "y": 2.5, "w": 2.2, "h": 2.0, "t": "3. WhatsApp\n\nEl votante envía mensaje pre-llenado al Chatbot oficial."},
        {"x": 9.8, "y": 2.5, "w": 2.7, "h": 2.0, "t": "4. CRM Petición\n\nSi tiene solicitud, el bot genera Folio GS-XXXX y envía al CRM."}
    ]

    for c in cajas_info:
        # Dibujar forma (rectángulo redondeado)
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c["x"]), Inches(c["y"]), Inches(c["w"]), Inches(c["h"]))
        aplicar_estilo_forma(rect, COLOR_GUINDA, COLOR_ORO, 1.5)
        
        # Agregar texto dentro del rectángulo
        tf_rect = rect.text_frame
        tf_rect.word_wrap = True
        p = tf_rect.paragraphs[0]
        p.text = c["t"]
        p.font.name = 'Arial'
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER

    # Flechas Conectoras
    flechas_info = [
        {"x": 3.1, "y": 3.3, "w": 0.6, "h": 0.4},
        {"x": 6.1, "y": 3.3, "w": 0.6, "h": 0.4},
        {"x": 9.1, "y": 3.3, "w": 0.6, "h": 0.4}
    ]
    for f in flechas_info:
        flecha = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(f["x"]), Inches(f["y"]), Inches(f["w"]), Inches(f["h"]))
        aplicar_estilo_forma(flecha, COLOR_ORO)

    # Texto al pie del diagrama
    notaBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0))
    tf_nota = notaBox.text_frame
    tf_nota.word_wrap = True
    p_nota = tf_nota.paragraphs[0]
    p_nota.text = "Nota: El flujo vincula el registro web de Firestore directamente con el Chatbot en WhatsApp para validar la línea celular del votante y, si lo requiere, transferirlo automáticamente al CRM de Vinculación."
    p_nota.font.name = 'Arial'
    p_nota.font.size = Pt(12)
    p_nota.font.italic = True
    p_nota.font.color.rgb = COLOR_TEXTO_OSCURO
    p_nota.alignment = PP_ALIGN.CENTER

    # ------------------------------------------
    # Diapositiva 6: Validación y Petición
    # ------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_FONDO)
    agregar_encabezado(slide, "Pasos 3, 4 y 5: Validación y Gestión")

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(7.2), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True

    bullets = [
        "El votante presiona 'Registrar' y se abre WhatsApp con un texto listo para enviar.",
        "Al enviar este mensaje, el Chatbot del Senador le da la bienvenida y valida sus datos.",
        "El Chatbot le preguntará al final si tiene alguna petición o necesita apoyo social.",
        "Si el votante describe su petición, el Chatbot genera de inmediato el Folio de Gestión GS-XXXX.",
        "El caso aparece al instante en el CRM administrativo (/vinculacion) para su seguimiento."
    ]
    for idx, b in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO_OSCURO
        p.space_before = Pt(12)

    # Tarjeta Destacada Lateral
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
    aplicar_estilo_forma(card, COLOR_GUINDA)
    
    txBox_card = slide.shapes.add_textbox(Inches(9.0), Inches(2.3), Inches(3.3), Inches(3.5))
    tf_card = txBox_card.text_frame
    tf_card.word_wrap = True
    p_card = tf_card.paragraphs[0]
    p_card.text = "SEGUIMIENTO DIGITAL\n\nEl folio de gestión GS-XXXX le permite al votante dar seguimiento a su caso escribiendo por WhatsApp o llamando a Vinculación (662 423 6390)."
    p_card.font.name = 'Arial'
    p_card.font.size = Pt(15)
    p_card.font.bold = True
    p_card.font.color.rgb = COLOR_BLANCO
    p_card.alignment = PP_ALIGN.CENTER

    # ------------------------------------------
    # Diapositiva 7: Cierre
    # ------------------------------------------
    slide = prs.slides.add_slide(slide_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    aplicar_estilo_forma(bg, COLOR_GUINDA)

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "¡Muchas Gracias por su Esfuerzo en Territorio!"
    p.font.name = 'Georgia'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORO
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Campana de Enlaces Ciudadanos de Heriberto Aguilar\nOrganización, Honestidad y Transformación para Sonora"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_BLANCO
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    # Guardar
    output_path = r"c:\Users\ENRIQ\SISTEMA DE VOTANTES PARA EL ESTADO DE SONORA\capacitacion_brigadistas.pptx"
    prs.save(output_path)
    print(f"Presentacion creada exitosamente en: {output_path}")

if __name__ == "__main__":
    crear_presentacion()
